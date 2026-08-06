# Copyright (c) 2025 Beijing Volcano Engine Technology Co., Ltd. and/or its affiliates.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Generate a PowerPoint deck and save it as an ADK artifact."""

from __future__ import annotations

import asyncio
import re
import tempfile
from pathlib import Path
from typing import Any

from google.adk.tools import ToolContext
from google.genai import types
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PPTX_MIME_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
PPT_PREVIEW_MIME_TYPE = "image/webp"
_MAX_SLIDES = 20
_MAX_BULLETS = 7
_MAX_SOURCES = 12

_THEMES: dict[str, dict[str, str]] = {
    "blue": {
        "canvas": "F4F7FB",
        "cover": "0B1F3A",
        "title": "10233F",
        "body": "344760",
        "muted": "6B7C93",
        "accent": "2F6FED",
        "cover_text": "FFFFFF",
        "cover_muted": "C9D8F0",
    },
    "dark": {
        "canvas": "10151E",
        "cover": "090D14",
        "title": "F5F7FA",
        "body": "CCD4E0",
        "muted": "8D9AAF",
        "accent": "74A7FF",
        "cover_text": "FFFFFF",
        "cover_muted": "AAB6C8",
    },
    "warm": {
        "canvas": "FBF6EF",
        "cover": "41281E",
        "title": "3D2A24",
        "body": "604A42",
        "muted": "8C7469",
        "accent": "D66A3A",
        "cover_text": "FFF8F2",
        "cover_muted": "E7CFC2",
    },
    "green": {
        "canvas": "F2F8F5",
        "cover": "12372B",
        "title": "173B31",
        "body": "36594F",
        "muted": "6D887F",
        "accent": "2B8A68",
        "cover_text": "F8FFFC",
        "cover_muted": "BFD9CF",
    },
}


def _clean_text(value: object, max_length: int) -> str:
    return " ".join(str(value or "").split())[:max_length]


def _string_list(value: object, *, limit: int, max_length: int) -> list[str]:
    if isinstance(value, str):
        values = value.splitlines()
    elif isinstance(value, list):
        values = value
    else:
        values = []
    return [text for item in values[:limit] if (text := _clean_text(item, max_length))]


def _normalize_slides(slides: list[dict[str, Any]]) -> list[dict[str, object]]:
    normalized: list[dict[str, object]] = []
    for index, slide in enumerate(slides[:_MAX_SLIDES]):
        if not isinstance(slide, dict):
            continue
        title = _clean_text(slide.get("title"), 120) or f"第 {index + 1} 页"
        normalized.append(
            {
                "title": title,
                "summary": _clean_text(slide.get("summary"), 240),
                "bullets": _string_list(
                    slide.get("bullets"), limit=_MAX_BULLETS, max_length=220
                ),
                "sources": _string_list(
                    slide.get("sources"), limit=_MAX_SOURCES, max_length=500
                ),
            }
        )
    return normalized


def _parse_deck_markdown(deck_markdown: str) -> list[dict[str, object]]:
    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for raw_line in deck_markdown.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("## "):
            if current is not None:
                slides.append(current)
            current = {
                "title": line.removeprefix("## ").strip(),
                "summary": "",
                "bullets": [],
                "sources": [],
            }
            continue
        if current is None:
            continue
        lowered = line.lower()
        if lowered.startswith("sources:") or line.startswith("来源："):
            value = line.split(":" if ":" in line else "：", 1)[1]
            current["sources"].extend(
                source.strip() for source in value.split("|") if source.strip()
            )
        elif line.startswith(("- ", "* ", "• ")):
            current["bullets"].append(line[2:].strip())
        elif not current["summary"]:
            current["summary"] = line
        else:
            current["bullets"].append(line)
    if current is not None:
        slides.append(current)
    return _normalize_slides(slides)


def _safe_filename(filename: str, title: str) -> str:
    value = Path(filename.strip()).name if filename.strip() else ""
    if not value:
        value = re.sub(r"[^\w\u4e00-\u9fff-]+", "-", title).strip("-")
    value = value[:100] or "presentation"
    if not value.lower().endswith(".pptx"):
        value += ".pptx"
    return value


async def _create_pptx(
    spec: dict[str, object],
    output_path: Path,
    preview_path: Path,
) -> None:
    await asyncio.to_thread(_write_presentation, spec, output_path, preview_path)
    if not output_path.is_file() or output_path.stat().st_size == 0:
        raise RuntimeError("PPT generation completed without producing a file.")
    if not preview_path.is_file() or preview_path.stat().st_size == 0:
        raise RuntimeError("PPT generation completed without producing a preview.")


def _color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _background(slide: Any, value: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _color(value)


def _text_box(
    slide: Any,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: str,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _color(color)


def _write_presentation(
    spec: dict[str, object],
    output_path: Path,
    preview_path: Path,
) -> None:
    theme_name = str(spec.get("theme") or "blue")
    theme = _THEMES.get(theme_name, _THEMES["blue"])
    raw_slides = spec.get("slides")
    slides = raw_slides if isinstance(raw_slides, list) else []

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    cover = presentation.slides.add_slide(blank_layout)
    _background(cover, theme["cover"])
    marker = cover.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.1), Inches(0.12), Inches(4.7)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = _color(theme["accent"])
    marker.line.fill.background()
    _text_box(
        cover,
        str(spec.get("title") or "演示文稿"),
        left=1.35,
        top=2.0,
        width=10.8,
        height=1.8,
        font_size=38,
        color=theme["cover_text"],
        bold=True,
    )
    subtitle = str(spec.get("subtitle") or "")
    if subtitle:
        _text_box(
            cover,
            subtitle,
            left=1.38,
            top=4.1,
            width=9.8,
            height=0.9,
            font_size=20,
            color=theme["cover_muted"],
        )

    for index, raw_slide in enumerate(slides):
        if not isinstance(raw_slide, dict):
            continue
        slide = presentation.slides.add_slide(blank_layout)
        _background(slide, theme["canvas"])
        _text_box(
            slide,
            str(raw_slide.get("title") or f"第 {index + 1} 页"),
            left=0.8,
            top=0.55,
            width=10.8,
            height=0.65,
            font_size=27,
            color=theme["title"],
            bold=True,
        )
        _text_box(
            slide,
            str(index + 2).zfill(2),
            left=11.7,
            top=0.62,
            width=0.7,
            height=0.4,
            font_size=12,
            color=theme["accent"],
            bold=True,
            alignment=PP_ALIGN.RIGHT,
        )
        summary = str(raw_slide.get("summary") or "")
        content_top = 1.55
        if summary:
            _text_box(
                slide,
                summary,
                left=0.85,
                top=content_top,
                width=11.35,
                height=0.72,
                font_size=18,
                color=theme["body"],
                bold=True,
            )
            content_top += 0.95
        bullets = raw_slide.get("bullets")
        bullet_values = bullets if isinstance(bullets, list) else []
        for bullet_index, bullet in enumerate(bullet_values[:_MAX_BULLETS]):
            _text_box(
                slide,
                f"•  {bullet}",
                left=1.0,
                top=content_top + bullet_index * 0.68,
                width=11.0,
                height=0.58,
                font_size=17 if len(bullet_values) >= 6 else 19,
                color=theme["body"],
            )
        sources = raw_slide.get("sources")
        source_values = sources if isinstance(sources, list) else []
        if source_values:
            _text_box(
                slide,
                "Sources: " + " | ".join(str(item) for item in source_values),
                left=0.85,
                top=6.9,
                width=11.6,
                height=0.3,
                font_size=8,
                color=theme["muted"],
            )

    presentation.save(output_path)
    _write_preview(spec, slides, theme, preview_path)


def _preview_font(
    size: int, *, bold: bool = False
) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    names = (
        "NotoSansCJK-Bold.ttc" if bold else "NotoSansCJK-Regular.ttc",
        "NotoSansCJK-Bold.otf" if bold else "NotoSansCJK-Regular.otf",
        "Hiragino Sans GB.ttc",
        "Arial Unicode.ttf",
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
    )
    roots = (
        Path("/usr/share/fonts/opentype/noto"),
        Path("/usr/share/fonts/truetype/noto"),
        Path("/usr/share/fonts/truetype/wqy"),
        Path("/System/Library/Fonts"),
        Path("/System/Library/Fonts/Supplemental"),
        Path("C:/Windows/Fonts"),
    )
    for root in roots:
        for name in names:
            candidate = root / name
            if candidate.is_file():
                try:
                    return ImageFont.truetype(str(candidate), size=size)
                except OSError:
                    continue
    return ImageFont.load_default(size=size)


def _fit_preview_text(text: object, limit: int) -> str:
    value = _clean_text(text, limit + 1)
    return value if len(value) <= limit else f"{value[: limit - 1]}…"


def _write_preview(
    spec: dict[str, object],
    slides: list[object],
    theme: dict[str, str],
    path: Path,
) -> None:
    width, height = 320, 180
    slide_count = len(slides) + 1
    columns = min(3, max(1, slide_count))
    rows = (slide_count + columns - 1) // columns
    montage = Image.new("RGB", (columns * width, rows * height), "#E7EBF1")
    draw = ImageDraw.Draw(montage)
    cover_title_font = _preview_font(19, bold=True)
    cover_subtitle_font = _preview_font(10)
    title_font = _preview_font(14, bold=True)
    summary_font = _preview_font(10, bold=True)
    body_font = _preview_font(10)
    page_font = _preview_font(9, bold=True)
    for index in range(slide_count):
        left = (index % columns) * width
        top = (index // columns) * height
        background = theme["cover"] if index == 0 else theme["canvas"]
        foreground = theme["cover_text"] if index == 0 else theme["title"]
        accent = theme["accent"]
        draw.rectangle(
            (left + 4, top + 4, left + width - 4, top + height - 4),
            fill=f"#{background}",
        )
        draw.rectangle(
            (left + 22, top + 28, left + 28, top + height - 28),
            fill=f"#{accent}",
        )
        if index == 0:
            draw.text(
                (left + 45, top + 54),
                _fit_preview_text(spec.get("title"), 22) or "演示文稿",
                fill=f"#{foreground}",
                font=cover_title_font,
            )
            subtitle = _fit_preview_text(spec.get("subtitle"), 34)
            if subtitle:
                draw.text(
                    (left + 46, top + 91),
                    subtitle,
                    fill=f"#{theme['cover_muted']}",
                    font=cover_subtitle_font,
                )
        else:
            raw_slide = slides[index - 1]
            slide = raw_slide if isinstance(raw_slide, dict) else {}
            draw.text(
                (left + 44, top + 34),
                _fit_preview_text(slide.get("title"), 30) or f"第 {index} 页",
                fill=f"#{foreground}",
                font=title_font,
            )
            content_top = top + 67
            summary = _fit_preview_text(slide.get("summary"), 42)
            if summary:
                draw.text(
                    (left + 45, content_top),
                    summary,
                    fill=f"#{theme['body']}",
                    font=summary_font,
                )
                content_top += 23
            bullets = slide.get("bullets")
            bullet_values = bullets if isinstance(bullets, list) else []
            for bullet_index, bullet in enumerate(bullet_values[:4]):
                draw.text(
                    (left + 46, content_top + bullet_index * 18),
                    f"• {_fit_preview_text(bullet, 38)}",
                    fill=f"#{theme['body']}",
                    font=body_font,
                )
        draw.text(
            (left + width - 33, top + 14),
            str(index + 1),
            fill=f"#{accent}",
            font=page_font,
        )
    montage.save(path, format="WEBP", quality=82, method=4)


async def ppt_generate(
    title: str,
    deck_markdown: str,
    tool_context: ToolContext,
    subtitle: str = "",
    theme: str = "blue",
    filename: str = "presentation.pptx",
) -> dict[str, object]:
    """Create a real PPTX file and attach it to the current conversation.

    Plan the complete deck before calling this tool. ``deck_markdown`` uses a
    simple flat format to avoid complex nested arguments: start every content
    slide with ``## Slide title``; add one plain-text summary line followed by
    3-7 ``- bullet`` lines. Add an optional ``Sources: URL | URL`` line when
    external claims or assets are used. The title slide is added automatically.

    Args:
        title: Audience-facing deck title.
        deck_markdown: Ordered content slides in the flat Markdown format.
        tool_context: Current ADK tool context used to save the result.
        subtitle: Optional title-slide subtitle.
        theme: Visual theme: blue, dark, warm, or green.
        filename: Download filename ending in .pptx.

    Returns:
        Metadata for the saved PowerPoint artifact.
    """
    deck_title = _clean_text(title, 160)
    if not deck_title:
        raise ValueError("title is required")
    normalized_slides = _parse_deck_markdown(deck_markdown)
    if not normalized_slides:
        raise ValueError(
            "deck_markdown must contain at least one content slide starting with ##"
        )
    artifact_name = _safe_filename(filename, deck_title)
    spec: dict[str, object] = {
        "title": deck_title,
        "subtitle": _clean_text(subtitle, 240),
        "theme": theme if theme in {"blue", "dark", "warm", "green"} else "blue",
        "slides": normalized_slides,
    }

    with tempfile.TemporaryDirectory(prefix="veadk-ppt-") as temp_dir:
        output_path = Path(temp_dir) / artifact_name
        preview_name = f"{Path(artifact_name).stem}.preview.webp"
        preview_path = Path(temp_dir) / preview_name
        await _create_pptx(spec, output_path, preview_path)
        preview_artifact = types.Part.from_bytes(
            data=preview_path.read_bytes(),
            mime_type=PPT_PREVIEW_MIME_TYPE,
        )
        if preview_artifact.inline_data is not None:
            preview_artifact.inline_data.display_name = preview_name
        preview_version = await tool_context.save_artifact(
            preview_name,
            preview_artifact,
        )
        artifact = types.Part.from_bytes(
            data=output_path.read_bytes(),
            mime_type=PPTX_MIME_TYPE,
        )
        if artifact.inline_data is not None:
            artifact.inline_data.display_name = artifact_name
        version = await tool_context.save_artifact(artifact_name, artifact)

    return {
        "status": "created",
        "filename": artifact_name,
        "version": version,
        "preview_filename": preview_name,
        "preview_version": preview_version,
        "slide_count": len(normalized_slides) + 1,
    }


__all__ = ["ppt_generate"]
